# /// script
# requires-python = ">=3.9"
# dependencies = ["python-pptx>=1.0.2"]
# ///
"""python-pptx layout builder for the google-slides-templating skill.

`LayoutShapes` (a slide layout's `.shapes`) supports almost no `add_*`
methods at all, not just pictures/connectors as previously documented,
confirmed empirically: its only public method is `clone_placeholder`.
Every new element (picture, autoshape, textbox, ...) has to be built on
a scratch slide first, where the full `SlideShapes` API is available,
then reparented into the layout's shape tree.
"""

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


def _reparent_to_layout(shape, layout):
    """Move a shape element from the scratch slide it was built on onto
    a layout's shape tree, re-registering any image relationship it
    depends on so the blip reference isn't left dangling."""
    element = shape._element
    source_part = shape.part

    for blip in element.findall(".//" + qn("a:blip")):
        r_embed = blip.get(qn("r:embed"))
        if r_embed:
            image_part = source_part.related_part(r_embed)
            new_rId = layout.part.relate_to(image_part, RT.IMAGE)
            blip.set(qn("r:embed"), new_rId)

    element.getparent().remove(element)
    layout.shapes._spTree.append(element)


def _remove_slide(prs, slide):
    """Delete a slide (no public python-pptx API for this): drop its
    sldId entry from the presentation's slide list, then its part/rel."""
    slide_id_lst = prs.slides._sldIdLst
    for sldId in list(slide_id_lst):
        if int(sldId.get("id")) == slide.slide_id:
            slide_id_lst.remove(sldId)
            break
    for rId, rel in list(prs.part.rels.items()):
        if rel.target_part is slide.part:
            prs.part.drop_rel(rId)
            break


def add_layout_picture(prs, layout, image_path, left, top, width=None, height=None):
    """Add a picture to a layout."""
    scratch = prs.slides.add_slide(prs.slide_layouts[6])
    picture = scratch.shapes.add_picture(image_path, left, top, width, height)
    _reparent_to_layout(picture, layout)
    _remove_slide(prs, scratch)
    return picture


def add_layout_shape(prs, layout, shape_type, left, top, width, height):
    """Add a plain autoshape to a layout."""
    scratch = prs.slides.add_slide(prs.slide_layouts[6])
    shape = scratch.shapes.add_shape(shape_type, left, top, width, height)
    _reparent_to_layout(shape, layout)
    _remove_slide(prs, scratch)
    return shape


def set_placeholder_position(layout, ph_idx, left, top, width, height):
    ph = layout.placeholders[ph_idx]
    ph.left, ph.top, ph.width, ph.height = left, top, width, height
    return ph


if __name__ == "__main__":
    import sys
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE

    out_path = sys.argv[1] if len(sys.argv) > 1 else "demo-layout.pptx"
    bg_path = sys.argv[2] if len(sys.argv) > 2 else "demo-bg.png"

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    add_layout_picture(
        prs, layout, bg_path, 0, 0, prs.slide_width, prs.slide_height
    )
    add_layout_shape(
        prs, layout, MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(1), Inches(1),
    )
    set_placeholder_position(
        layout, 0, Inches(1), Inches(0.5), Inches(8), Inches(1)
    )

    prs.slides.add_slide(layout)
    prs.save(out_path)
    print(f"Wrote {out_path}")
